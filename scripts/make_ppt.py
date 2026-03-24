from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colors ──────────────────────────────────────────────
RED    = RGBColor(0xE6, 0x00, 0x12)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
MID    = RGBColor(0x66, 0x66, 0x66)
LIGHT  = RGBColor(0xF8, 0xF9, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xFF, 0x95, 0x00)
BLUE   = RGBColor(0x00, 0x7A, 0xFF)
BORDER = RGBColor(0xEE, 0xEE, 0xEE)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── 헬퍼 ──────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = Pt(border_pt)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def ellipse(slide, x, y, w, h, fill_color=RED):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    prstGeom = shape._element.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'ellipse')
    return shape

def step_badge(slide, num, x, y, col=RED):
    c = ellipse(slide, x, y, 0.46, 0.46, fill_color=col)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = WHITE

def header(slide, title, col=RED):
    box(slide, 0, 0, 13.33, 1.3, fill_color=col)
    txt(slide, title, 0.6, 0.25, 12, 0.8, size=32, bold=True, color=WHITE)

def who_badge(slide, label, col, x, y):
    box(slide, x, y, 2.1, 0.36, fill_color=col)
    txt(slide, label, x+0.1, y+0.04, 1.9, 0.28, size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 타이틀
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, DARK)
box(slide, 0, 0, 0.55, 7.5, fill_color=RED)

txt(slide, "Lego Lineup", 1.0, 1.6, 10, 1.1,
    size=56, bold=True, italic=True, color=RED)
txt(slide, "강사 사용 매뉴얼", 1.0, 2.8, 10, 0.9,
    size=38, bold=True, color=WHITE)
txt(slide, "실시간 레고 팀 활동 관리 도구  |  theplaycompany", 1.05, 3.82, 10, 0.5,
    size=16, color=RGBColor(0x88, 0x88, 0x88))

box(slide, 1.0, 5.6, 11.7, 0.04, fill_color=RGBColor(0x44, 0x44, 0x44))
txt(slide, "관리자 주소:  /admin", 1.05, 5.8, 6, 0.4,
    size=14, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA))


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 전체 진행 순서 한눈에 보기
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "전체 진행 순서", DARK)

txt(slide, "5단계 흐름  |  화살표 순서대로 진행하세요", 0.6, 1.42, 12, 0.4,
    size=14, color=MID)

steps = [
    ("1", "기록 입력",    "강사 (다수)",  RED,    "팀별 수행 시간\n입력"),
    ("2", "색상 순서\n작성", "교육생",   BLUE,   "5가지 색상\n순서 제출"),
    ("3", "턴수 활성화\n버튼 클릭", "강사 1명", GREEN, "/admin에서\n버튼 클릭"),
    ("4", "턴수 작성",    "교육생",      BLUE,   "전략 턴수\n입력 & 제출"),
    ("5", "분석 버튼\n클릭", "강사 1명",  ORANGE, "통계 탭\n분석 확인"),
]

card_w = 2.25
gap    = 0.22
start  = (13.33 - (card_w * 5 + gap * 4)) / 2

for i, (num, title, who, col, desc) in enumerate(steps):
    cx = start + i * (card_w + gap)

    # 카드 배경
    box(slide, cx, 2.0, card_w, 4.6,
        fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=col, border_pt=2)

    # 번호 뱃지
    step_badge(slide, num, cx + card_w/2 - 0.23, 2.15, col=col)

    # 누가 하는지 태그
    tag_col = RED if "강사" in who else BLUE
    box(slide, cx + 0.15, 2.78, card_w - 0.3, 0.33, fill_color=tag_col)
    txt(slide, who, cx + 0.2, 2.82, card_w - 0.35, 0.26,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 제목
    txt(slide, title, cx + 0.1, 3.22, card_w - 0.2, 0.8,
        size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # 설명
    txt(slide, desc, cx + 0.1, 4.12, card_w - 0.2, 0.8,
        size=12, color=MID, align=PP_ALIGN.CENTER)

    # 화살표 (마지막 제외)
    if i < len(steps) - 1:
        ax = cx + card_w + gap/2 - 0.08
        txt(slide, "→", ax, 3.9, 0.3, 0.4, size=20, bold=True, color=BORDER)

# 범례
box(slide, 0.5, 6.8, 12.3, 0.45, fill_color=LIGHT, border_color=BORDER, border_pt=1)
box(slide, 0.7, 6.9, 0.7, 0.26, fill_color=RED)
txt(slide, "강사", 1.45, 6.9, 1.2, 0.26, size=11, color=DARK)
box(slide, 2.7, 6.9, 0.7, 0.26, fill_color=BLUE)
txt(slide, "교육생", 3.45, 6.9, 1.5, 0.26, size=11, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — STEP 1: 기록 입력 (강사)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "STEP 1  —  기록 입력", RED)

who_badge(slide, "강사 (10~15명)", RED, 0.6, 1.42)
txt(slide, "교육 활동 완료 후, 각 강사가 담당 팀의 수행 시간을 입력합니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

# 좌: 입력 방법
box(slide, 0.5, 2.0, 6.0, 4.9, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "입력 방법", 0.8, 2.15, 5, 0.4, size=17, bold=True, color=DARK)

row_steps = [
    ("①", "/admin 접속 후 [기록 입력] 탭 클릭"),
    ("②", "팀 ID 번호 입력"),
    ("③", "분(1~20) 선택  +  초(0~59) 선택"),
    ("④", "[기록 저장] 버튼 클릭"),
    ("⑤", "다른 팀도 동일하게 반복"),
]
for i, (num, step) in enumerate(row_steps):
    by = 2.72 + i * 0.72
    txt(slide, num, 0.75, by, 0.45, 0.45, size=16, bold=True, color=RED)
    txt(slide, step, 1.3, by, 5.0, 0.5, size=14, color=DARK)

# 우: 주의사항 + 팁
box(slide, 6.9, 2.0, 6.0, 2.2, fill_color=RGBColor(0xFF,0xF0,0xF0),
    border_color=RED, border_pt=1)
txt(slide, "⚠  주의사항", 7.15, 2.12, 5.5, 0.38, size=14, bold=True, color=RED)
txt(slide,
    "• 팀당 최대 3개 기록 입력 가능\n"
    "• 잘못 입력 시 [수정] 또는 [삭제] 후 재입력\n"
    "• 기록 형식 예:  3' 24\"",
    7.15, 2.6, 5.5, 1.3, size=13, color=DARK)

box(slide, 6.9, 4.4, 6.0, 2.5, fill_color=RGBColor(0xF0,0xFF,0xF4),
    border_color=GREEN, border_pt=1)
txt(slide, "💡  팁", 7.15, 4.52, 5.5, 0.38, size=14, bold=True, color=GREEN)
txt(slide,
    "• 강사 여러 명이 동시에 입력 가능\n"
    "• 실시간으로 반영되므로 중복 입력 주의\n"
    "• 3팀 이상 입력되면 통계 탭에서 분석 가능",
    7.15, 4.95, 5.5, 1.7, size=13, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — STEP 2: 색상 순서 작성 (교육생)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "STEP 2  —  색상 순서 작성", BLUE)

who_badge(slide, "교육생 (각 팀)", BLUE, 0.6, 1.42)
txt(slide, "교육생이 팀별로 URL에 접속하여 색상 순서를 제출합니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

# 교육생 화면 단계
flow = [
    ("①", "QR 또는 URL로 접속"),
    ("②", "팀 ID 번호 입력"),
    ("③", "5가지 색상을 원하는 순서로 선택\n(빨강 · 노랑 · 파랑 · 초록 · 흰색)"),
    ("④", "[제출하기] 버튼 클릭 → 대기 화면으로 이동"),
]
for i, (num, step) in enumerate(flow):
    cy = 2.1 + i * 1.05
    step_badge(slide, num, 0.6, cy, col=BLUE)
    box(slide, 1.25, cy, 5.5, 0.88,
        fill_color=LIGHT, border_color=BORDER, border_pt=1)
    txt(slide, step, 1.45, cy+0.18, 5.2, 0.65, size=14, color=DARK)

# 우: 강사 확인 포인트
box(slide, 7.2, 2.0, 5.7, 4.9, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "강사 확인 포인트", 7.5, 2.15, 5.0, 0.4, size=17, bold=True, color=DARK)
txt(slide, "/admin 대시보드 탭에서 확인", 7.5, 2.6, 5.0, 0.35,
    size=12, italic=True, color=MID)

checks = [
    "테이블의 [색상] 열에 색상 원이 표시되면 제출 완료",
    "색상 열이 비어 있는 팀 ID = 미제출 팀",
    "전원 제출 완료 확인 후 STEP 3으로 진행",
]
for i, c in enumerate(checks):
    cy = 3.1 + i * 0.88
    txt(slide, "✓", 7.4, cy, 0.4, 0.45, size=16, bold=True, color=GREEN)
    txt(slide, c, 7.85, cy, 4.8, 0.7, size=13, color=DARK)

box(slide, 7.3, 5.6, 5.5, 1.1, fill_color=RGBColor(0xFF,0xF7,0xE6),
    border_color=ORANGE, border_pt=1)
txt(slide, "⚠  색상 잘못 제출 시", 7.5, 5.7, 5.0, 0.35, size=12, bold=True, color=ORANGE)
txt(slide, "참가자 화면에서 [수정하기] 버튼으로 재입력 가능",
    7.5, 6.05, 5.0, 0.5, size=12, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — STEP 3: 턴수 활성화 버튼 클릭 (강사 1명)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "STEP 3  —  턴수 활성화 버튼 클릭", GREEN)

who_badge(slide, "강사 1명", GREEN, 0.6, 1.42)
txt(slide, "모든 팀이 색상 순서를 제출한 것을 확인한 후, 강사 1명이 버튼을 클릭합니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

# 좌: 절차
box(slide, 0.5, 2.05, 6.0, 4.8, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "진행 절차", 0.8, 2.2, 5, 0.4, size=17, bold=True, color=DARK)

proc = [
    ("①", "전체 색상 제출 완료 확인 (대시보드 탭)"),
    ("②", "/admin 상단의 [턴 수 활성화] 버튼 클릭"),
    ("③", "확인 메시지 → [확인] 클릭"),
    ("④", "버튼이 초록색으로 변하면 활성화 완료"),
    ("⑤", "참가자 대기 화면 → [다음 단계] 버튼 해제"),
]
for i, (num, step) in enumerate(proc):
    by = 2.78 + i * 0.72
    txt(slide, num, 0.75, by, 0.45, 0.45, size=16, bold=True, color=GREEN)
    txt(slide, step, 1.3, by, 5.0, 0.5, size=14, color=DARK)

# 우 상단: 버튼 상태
box(slide, 6.9, 2.05, 6.0, 2.0, fill_color=RGBColor(0xF0,0xFF,0xF4),
    border_color=GREEN, border_pt=2)
txt(slide, "버튼 상태 의미", 7.15, 2.18, 5.5, 0.38, size=14, bold=True, color=DARK)
box(slide, 7.1, 2.68, 2.4, 0.5, fill_color=RGBColor(0xEE,0xEE,0xEE),
    border_color=BORDER, border_pt=1)
txt(slide, "턴 수 활성화", 7.2, 2.75, 2.2, 0.35, size=13, bold=True,
    color=RED, align=PP_ALIGN.CENTER)
txt(slide, "← 비활성 (빨간 테두리)", 9.6, 2.75, 3.2, 0.35, size=12, color=MID)
box(slide, 7.1, 3.28, 2.4, 0.5, fill_color=GREEN)
txt(slide, "턴 수 중단", 7.2, 3.35, 2.2, 0.35, size=13, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "← 활성화 완료 (초록)", 9.6, 3.35, 3.2, 0.35, size=12, color=MID)

# 우 하단: 주의
box(slide, 6.9, 4.25, 6.0, 2.6, fill_color=RGBColor(0xFF,0xF0,0xF0),
    border_color=RED, border_pt=1)
txt(slide, "⚠  주의사항", 7.15, 4.38, 5.5, 0.38, size=14, bold=True, color=RED)
txt(slide,
    "• 강사 1명만 클릭 (중복 클릭 불필요)\n"
    "• 실수로 껐다가 다시 켜도 기존 데이터에 영향 없음\n"
    "• 활성화 전 이미 대기 화면인 팀은 즉시 다음 단계로 이동",
    7.15, 4.85, 5.5, 1.7, size=13, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — STEP 4: 턴수 작성 (교육생)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "STEP 4  —  턴수 작성", BLUE)

who_badge(slide, "교육생 (각 팀)", BLUE, 0.6, 1.42)
txt(slide, "강사가 턴수를 활성화하면 교육생 화면에 [다음 단계 진행하기] 버튼이 나타납니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

flow = [
    ("①", "[다음 단계 진행하기] 버튼 클릭"),
    ("②", "앞 드롭다운: 앞 숫자 T  선택 (2T ~ 5T)"),
    ("③", "뒤 드롭다운: + 뒤 숫자  선택 (+1 ~ +4)"),
    ("④", "[최종 제출] 버튼 클릭 → 완료 화면"),
]
for i, (num, step) in enumerate(flow):
    cy = 2.1 + i * 1.05
    step_badge(slide, num, 0.6, cy, col=BLUE)
    box(slide, 1.25, cy, 5.5, 0.88,
        fill_color=LIGHT, border_color=BORDER, border_pt=1)
    txt(slide, step, 1.45, cy+0.18, 5.2, 0.65, size=14, color=DARK)

# 예시 박스
box(slide, 0.5, 6.25, 6.25, 0.9, fill_color=RGBColor(0xFF,0xF7,0xE6),
    border_color=ORANGE, border_pt=1)
txt(slide, "입력 예시:  3T + 2  →  \"3T+2\" 로 저장됨",
    0.75, 6.42, 5.8, 0.45, size=13, color=DARK)

# 우: 강사 확인
box(slide, 7.2, 2.0, 5.7, 4.9, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "강사 확인 포인트", 7.5, 2.15, 5.0, 0.4, size=17, bold=True, color=DARK)
txt(slide, "/admin 대시보드 탭에서 확인", 7.5, 2.6, 5.0, 0.35,
    size=12, italic=True, color=MID)

checks = [
    "테이블 [턴 수] 열에 값이 표시되면 제출 완료",
    "턴 수 열이 비어 있는 팀 ID = 미제출 팀",
    "전원 제출 완료 후 STEP 5로 진행",
]
for i, c in enumerate(checks):
    cy = 3.1 + i * 0.88
    txt(slide, "✓", 7.4, cy, 0.4, 0.45, size=16, bold=True, color=GREEN)
    txt(slide, c, 7.85, cy, 4.8, 0.7, size=13, color=DARK)

box(slide, 7.3, 5.6, 5.5, 1.1, fill_color=RGBColor(0xFF,0xF0,0xF0),
    border_color=RED, border_pt=1)
txt(slide, "⚠  턴수 미제출 팀 발견 시", 7.5, 5.7, 5.0, 0.35, size=12, bold=True, color=RED)
txt(slide, "팀 ID 확인 후 해당 팀에게 직접 안내",
    7.5, 6.05, 5.0, 0.5, size=12, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — STEP 5: 분석 버튼 클릭 (강사 1명)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "STEP 5  —  통계 탭 분석 확인", ORANGE)

who_badge(slide, "강사 1명", ORANGE, 0.6, 1.42)
txt(slide, "기록이 충분히 입력되면 통계 탭에서 분석 인사이트를 확인합니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

# 접근 방법
box(slide, 0.5, 2.0, 5.5, 1.5, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "접근 방법", 0.8, 2.12, 5.0, 0.38, size=15, bold=True, color=DARK)
txt(slide, "/admin 하단 탭 →  [통계]  →  [분석 보기] 버튼 클릭",
    0.8, 2.6, 5.0, 0.7, size=13, color=DARK)

# 인사이트 3개
insights = [
    (ORANGE, "인사이트 1", "기록 상위 팀의 전략 패턴",
     "상위 20% 팀의 색상 순서 패턴과\n평균 턴 수를 자동 분석"),
    (RED,    "인사이트 2", "동일 턴 수 내 기록 차이",
     "같은 턴 수인데 시간 차이가 30초\n이상인 경우 자동 감지"),
    (BLUE,   "인사이트 3", "턴 수별 평균 기록 시각화",
     "2T~5T 구간별 평균 수행 시간\n막대 그래프 자동 생성"),
]
for i, (col, label, title, desc) in enumerate(insights):
    cy = 3.7 + i * 1.15
    box(slide, 0.5, cy, 5.5, 1.05, fill_color=WHITE, border_color=col, border_pt=2)
    txt(slide, label, 0.72, cy+0.1, 1.8, 0.3, size=11, bold=True, color=col)
    txt(slide, title, 0.72, cy+0.38, 5.0, 0.35, size=14, bold=True, color=DARK)
    txt(slide, desc, 0.72, cy+0.68, 5.0, 0.3, size=12, color=MID)

# 우: 디브리핑 활용
box(slide, 6.4, 2.0, 6.5, 4.9, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "디브리핑 활용 포인트", 6.7, 2.15, 6.0, 0.4, size=17, bold=True, color=DARK)

debriefs = [
    ("Q1", "기록이 빠른 팀은 색상 순서가 어떻게 달랐나요?"),
    ("Q2", "같은 턴 수인데 시간 차이가 나는 이유는?"),
    ("Q3", "더 적은 턴으로 빠르게 완성한 팀의 전략은?"),
    ("Q4", "다시 시도한다면 어떻게 전략을 바꾸겠나요?"),
]
for i, (q, content) in enumerate(debriefs):
    cy = 2.75 + i * 1.0
    box(slide, 6.55, cy, 0.55, 0.55, fill_color=ORANGE)
    tf = slide.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = q
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE
    txt(slide, content, 7.2, cy+0.08, 5.5, 0.5, size=13, color=DARK)

box(slide, 0.5, 7.05, 12.3, 0.28, fill_color=RGBColor(0xF0,0xFF,0xF4),
    border_color=GREEN, border_pt=1)
txt(slide, "※ 3팀 이상 기록이 입력되어야 분석 인사이트가 표시됩니다.",
    0.7, 7.08, 12.0, 0.22, size=11, italic=True, color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — 세션 종료 & 초기화
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
bg(slide, WHITE)
header(slide, "세션 종료  &  데이터 초기화", DARK)

who_badge(slide, "강사 1명", DARK, 0.6, 1.42)
txt(slide, "교육 세션이 완전히 종료된 후, 다음 세션을 위해 데이터를 초기화합니다.",
    2.85, 1.46, 10, 0.36, size=13, color=MID)

# 초기화 절차
box(slide, 0.5, 2.05, 6.0, 4.2, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "초기화 절차", 0.8, 2.2, 5, 0.4, size=17, bold=True, color=DARK)

init_steps = [
    ("①", "/admin 대시보드 탭으로 이동"),
    ("②", "[데이터 초기화] 버튼 클릭"),
    ("③", "확인 팝업에서 '초기화' 직접 입력"),
    ("④", "[초기화] 버튼 클릭 → 전체 삭제 완료"),
]
for i, (num, step) in enumerate(init_steps):
    by = 2.85 + i * 0.77
    txt(slide, num, 0.75, by, 0.45, 0.45, size=16, bold=True, color=RED)
    txt(slide, step, 1.3, by, 5.0, 0.5, size=14, color=DARK)

# 경고
box(slide, 0.5, 5.55, 6.0, 1.3, fill_color=RGBColor(0xFF,0xF0,0xF0),
    border_color=RED, border_pt=2)
txt(slide, "⚠  삭제된 데이터는 복구 불가합니다.", 0.75, 5.65, 5.5, 0.4,
    size=13, bold=True, color=RED)
txt(slide, "반드시 세션 완전 종료 후 진행하세요.",
    0.75, 6.05, 5.5, 0.65, size=13, color=DARK)

# 우: 체크리스트
box(slide, 6.9, 2.05, 6.0, 4.8, fill_color=LIGHT, border_color=BORDER, border_pt=1)
txt(slide, "세션 종료 체크리스트", 7.2, 2.2, 5.5, 0.4, size=17, bold=True, color=DARK)

checks = [
    "분석 결과 화면 캡처 완료",
    "디브리핑 완료",
    "턴수 버튼 비활성화 상태 확인",
    "모든 교육생 완료 화면 확인",
    "데이터 초기화 완료",
]
for i, c in enumerate(checks):
    cy = 2.85 + i * 0.77
    box(slide, 7.1, cy+0.05, 0.34, 0.34, fill_color=WHITE,
        border_color=BORDER, border_pt=1)
    txt(slide, c, 7.55, cy, 5.1, 0.5, size=14, color=DARK)


# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
out = "/Users/stonylee/workspace/legolineup/LegoLineup_강사매뉴얼.pptx"
prs.save(out)
print(f"저장 완료: {out}")
